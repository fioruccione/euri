"""
Estrazione testo da formati di documento supportati, con cascata
testo-nativo → Vision LLM (Gemma 4 multimodale) per i casi scansionati.

Formati gestiti:
- PDF: pypdf → (fallback Vision per scansionati)
- DOCX (Word): python-docx (sempre testuale, niente fallback)
- PPTX (PowerPoint): python-pptx → (fallback Vision per slide grafiche pure)
- Immagini (jpg/png/etc): direttamente Vision (no fallback, è il primario)

Idea uniforme: l'estrattore restituisce sempre una STRINGA con il testo
del documento, formattato con marker `[Pagina N]` / `[Slide N]`. La
stringa viene poi iniettata nel prompt di Gemma per la generazione
del codice di analisi (CodeRunner pipeline).
"""
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


# Soglia minima di caratteri estratti da pypdf prima di considerare
# il PDF "scansionato" e attivare il fallback Vision.
EMPTY_PDF_THRESHOLD = 50


def extract_pdf_text(
    path: str | Path,
    vision_callback=None,
    dpi: int = 200,
) -> str:
    """
    Estrae il testo da un PDF con cascata pypdf → Vision.

    Args:
        path: percorso al PDF.
        vision_callback: funzione opzionale `(image_path: str, question: str) -> str`
            chiamata per ogni pagina se pypdf non estrae nulla.
            Tipicamente Brain.analyze_image (Gemma 4 multimodale).
            Se None, il fallback Vision è disabilitato e per i PDF
            scansionati si restituisce stringa vuota.
        dpi: risoluzione delle PNG per il fallback Vision. 200 è un buon
            compromesso lettura/velocità per A4 stampate.

    Returns:
        Testo estratto, con marker `[Pagina N]` davanti a ciascuna pagina.
        Stringa vuota se l'estrazione fallisce totalmente o non c'è
        callback Vision per un PDF scansionato.
    """
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(f"PDF non trovato: {path}")

    # Cascata 1: pypdf (PDF testuali — il caso più comune)
    pypdf_pages = _extract_with_pypdf(path)
    pypdf_total_chars = sum(len(p) for p in pypdf_pages)

    if pypdf_total_chars >= EMPTY_PDF_THRESHOLD:
        return _format_pages(pypdf_pages)

    # Cascata 2: Vision LLM via callback (PDF scansionati)
    if vision_callback is None:
        logger.warning(
            f"PDF '{path.name}' sembra scansionato ({pypdf_total_chars} char "
            f"estratti da pypdf) ma non c'è vision_callback — testo non disponibile"
        )
        return ""

    logger.info(
        f"PDF '{path.name}': testo pypdf insufficiente "
        f"({pypdf_total_chars} char) → fallback Vision"
    )
    vision_pages = _extract_with_vision(path, vision_callback, dpi)
    return _format_pages(vision_pages)


def _extract_with_pypdf(path: Path) -> list[str]:
    """Estrazione testo nativa via pypdf. Restituisce lista per pagina."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        return [(page.extract_text() or "").strip() for page in reader.pages]
    except Exception as e:
        logger.debug(f"pypdf fallito su {path.name}: {e}")
        return []


def _extract_with_vision(path: Path, vision_callback, dpi: int) -> list[str]:
    """
    Converte ciascuna pagina del PDF in PNG e chiede al vision_callback
    di estrarre il testo. Restituisce lista per pagina.
    """
    try:
        from pdf2image import convert_from_path
    except ImportError as e:
        logger.error(f"pdf2image non disponibile: {e}")
        return []

    try:
        images = convert_from_path(str(path), dpi=dpi)
    except Exception as e:
        logger.error(f"Conversione PDF→PNG fallita per {path.name}: {e}")
        return []

    # Prompt mirato a estrazione, non descrizione discorsiva
    question = (
        "Questa è una pagina di un documento tecnico, amministrativo o "
        "industriale (scheda tecnica, certificato, report, etichetta). "
        "Estrai TUTTO il testo che vedi, mantenendo l'organizzazione "
        "originale: tabelle come tabelle (righe separate da newline, "
        "colonne separate da | quando possibile), titoli su righe a sé, "
        "valori numerici con le loro unità di misura. Non commentare, "
        "non riassumere, non spiegare: riporta solo il contenuto integrale "
        "del documento."
    )

    pages_text = []
    for i, img in enumerate(images, 1):
        # Salvo l'immagine in un file temporaneo (analyze_image vuole un path)
        tmp_path = Path(f"/tmp/euri_pdf_page_{path.stem}_{i}.png")
        try:
            img.save(str(tmp_path), format="PNG")
            text = vision_callback(str(tmp_path), question) or ""
            pages_text.append(text.strip())
        except Exception as e:
            logger.error(f"Vision fallito su pagina {i} di {path.name}: {e}")
            pages_text.append("")
        finally:
            try:
                tmp_path.unlink(missing_ok=True)
            except Exception:
                pass

    return pages_text


def _format_pages(pages: list[str], marker: str = "Pagina") -> str:
    """Concatena le pagine con marker, ignorando quelle vuote."""
    blocks = []
    for i, text in enumerate(pages, 1):
        if text:
            blocks.append(f"[{marker} {i}]\n{text}")
    return "\n\n".join(blocks).strip()


# ──────────────────────────────────────────
# WORD .docx (Microsoft Word moderno)
# ──────────────────────────────────────────

def extract_docx_text(path: str | Path) -> str:
    """
    Estrae il testo da un documento Word .docx via python-docx.
    Include paragrafi e celle di tabelle. Niente fallback Vision: i .docx
    sono per definizione strutturati (XML zippato), se sono scansionati
    sono in realtà PDF rinominati o immagini incollate.
    """
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(f"DOCX non trovato: {path}")

    try:
        from docx import Document
    except ImportError as e:
        logger.error(f"python-docx non installato: {e}")
        return ""

    try:
        doc = Document(str(path))
        blocks = []

        # Paragrafi del corpo
        for p in doc.paragraphs:
            t = (p.text or "").strip()
            if t:
                blocks.append(t)

        # Tabelle: riga per riga, celle separate da " | "
        for i, table in enumerate(doc.tables, 1):
            blocks.append(f"\n[Tabella {i}]")
            for row in table.rows:
                cells = [(c.text or "").strip() for c in row.cells]
                if any(cells):
                    blocks.append(" | ".join(cells))

        return "\n".join(blocks).strip()
    except Exception as e:
        logger.error(f"python-docx fallito su {path.name}: {e}")
        return ""


# ──────────────────────────────────────────
# POWERPOINT .pptx
# ──────────────────────────────────────────

def extract_pptx_text(
    path: str | Path,
    vision_callback=None,
    dpi: int = 150,
) -> str:
    """
    Estrae il testo da una presentazione PowerPoint .pptx.

    Cascata:
      1. python-pptx: testo nativo da shapes/textframes (slide testuali, 95%
         dei casi commerciali — slide con bullet, titoli, tabelle)
      2. Fallback Vision se il contenuto testuale è troppo scarso: rasterizza
         ogni slide in PNG e chiede a Gemma 4 di leggerla. Utile per slide
         interamente grafiche o immagini incollate senza testo selezionabile.

    Args:
        path: percorso al .pptx
        vision_callback: come per PDF (es. Brain.analyze_image)
        dpi: risoluzione PNG per fallback (default 150 — slide ~16:9 a 1280×720)

    Returns:
        Testo estratto con marker `[Slide N]`.
    """
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX non trovato: {path}")

    pptx_slides = _extract_pptx_native(path)
    pptx_total_chars = sum(len(s) for s in pptx_slides)

    if pptx_total_chars >= EMPTY_PDF_THRESHOLD:
        return _format_pages(pptx_slides, marker="Slide")

    if vision_callback is None:
        logger.warning(
            f"PPTX '{path.name}' sembra grafico ({pptx_total_chars} char "
            f"estratti) ma non c'è vision_callback — testo non disponibile"
        )
        return ""

    logger.info(
        f"PPTX '{path.name}': testo python-pptx insufficiente "
        f"({pptx_total_chars} char) → fallback Vision"
    )
    vision_slides = _extract_pptx_via_vision(path, vision_callback, dpi)
    return _format_pages(vision_slides, marker="Slide")


def _extract_pptx_native(path: Path) -> list[str]:
    """Estrae testo da shapes e tabelle di ogni slide via python-pptx."""
    try:
        from pptx import Presentation
    except ImportError as e:
        logger.error(f"python-pptx non installato: {e}")
        return []

    try:
        prs = Presentation(str(path))
    except Exception as e:
        logger.error(f"python-pptx fallito su {path.name}: {e}")
        return []

    slides_text = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            # Testo da text_frame
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        parts.append(t)
            # Tabelle
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [(c.text or "").strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        slides_text.append("\n".join(parts).strip())
    return slides_text


def _extract_pptx_via_vision(path: Path, vision_callback, dpi: int) -> list[str]:
    """
    Fallback Vision: converte slide in PNG via LibreOffice (headless) e
    le passa al vision_callback. Richiede `libreoffice` installato.

    NOTA: la conversione pptx→png con LibreOffice è meno diretta di
    pdf→png. Se LibreOffice non è disponibile, restituisce lista vuota.
    """
    import subprocess
    import tempfile

    question = (
        "Questa è una slide di una presentazione. Estrai TUTTO il testo "
        "che vedi (titoli, bullet, didascalie, testo dentro grafici/immagini). "
        "Mantieni l'organizzazione gerarchica. Non commentare."
    )

    try:
        # LibreOffice converte pptx→pdf in dir temp, poi pdf2image per le PNG
        with tempfile.TemporaryDirectory() as tmpdir:
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf",
                 "--outdir", tmpdir, str(path)],
                check=True, capture_output=True, timeout=60,
            )
            pdf_path = Path(tmpdir) / (path.stem + ".pdf")
            if not pdf_path.exists():
                logger.error(f"LibreOffice conversion non ha prodotto PDF per {path.name}")
                return []

            from pdf2image import convert_from_path
            images = convert_from_path(str(pdf_path), dpi=dpi)
    except FileNotFoundError:
        logger.error("LibreOffice non installato — impossibile fallback Vision su PPTX")
        return []
    except Exception as e:
        logger.error(f"Conversione PPTX→PNG fallita per {path.name}: {e}")
        return []

    slides_text = []
    for i, img in enumerate(images, 1):
        tmp_path = Path(f"/tmp/euri_pptx_slide_{path.stem}_{i}.png")
        try:
            img.save(str(tmp_path), format="PNG")
            text = vision_callback(str(tmp_path), question) or ""
            slides_text.append(text.strip())
        except Exception as e:
            logger.error(f"Vision fallito su slide {i} di {path.name}: {e}")
            slides_text.append("")
        finally:
            try:
                tmp_path.unlink(missing_ok=True)
            except Exception:
                pass

    return slides_text


# ──────────────────────────────────────────
# IMMAGINI (jpg/png/webp/etc)
# ──────────────────────────────────────────

# Estensioni che riconosciamo come immagini estraibili via Vision
IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".tiff", ".gif"}


def extract_image_description(
    path: str | Path,
    vision_callback,
) -> str:
    """
    Wrapper coerente con gli altri extractor: passa l'immagine a Gemma 4
    multimodale e ottiene una descrizione/lettura del contenuto.

    A differenza di pdf/pptx, qui non c'è una "via testuale" da provare
    prima — Vision è il primo e unico canale. Se il vision_callback non
    è disponibile, ritorna stringa vuota.

    NOTA: il prompt è più ricco di quello PDF perché un'immagine singola
    può essere foto/screenshot/etichetta/documento — chiediamo a Gemma
    di adattare la risposta al tipo che vede.
    """
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(f"Immagine non trovata: {path}")
    if vision_callback is None:
        return ""

    question = (
        "Questa è un'immagine: può essere una foto, uno screenshot, un'etichetta, "
        "un certificato, una scheda tecnica, un grafico, o altro. "
        "1) Identifica il tipo di immagine in una frase. "
        "2) Estrai TUTTO il testo che vedi (se ce n'è), mantenendo "
        "l'organizzazione (tabelle come tabelle, titoli su righe a sé). "
        "3) Se è una foto senza testo, descrivi cosa rappresenta in italiano "
        "in 2-3 frasi concrete. "
        "Non commentare, non spiegare il tuo ragionamento."
    )

    try:
        text = vision_callback(str(path), question) or ""
        return text.strip()
    except Exception as e:
        logger.error(f"Vision fallito su immagine {path.name}: {e}")
        return ""


# ──────────────────────────────────────────
# DISPATCHER per estensione
# ──────────────────────────────────────────

def extract_any(path: str | Path, vision_callback=None) -> str:
    """
    Dispatcher generico: identifica il formato dall'estensione e chiama
    l'estrattore appropriato. Ritorna stringa vuota per formati non
    supportati (txt/md/json/csv NON vengono pre-estratti — Gemma li
    leggerà direttamente via pandas/open() nel codice generato).
    """
    path = Path(path)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return extract_pdf_text(path, vision_callback=vision_callback)
    if suffix == ".docx":
        return extract_docx_text(path)
    if suffix == ".pptx":
        return extract_pptx_text(path, vision_callback=vision_callback)
    if suffix in IMAGE_EXTS:
        return extract_image_description(path, vision_callback=vision_callback)

    # csv/xlsx/json/txt/md: niente pre-estrazione, sono leggibili da Gemma via codice
    return ""
